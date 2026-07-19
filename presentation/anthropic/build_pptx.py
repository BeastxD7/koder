#!/usr/bin/env python3
"""Builds lakshx-pitch-deck.pptx from the same content/palette as
lakshx-pitch-deck.html. Regenerate after editing either deck by re-running
this script (`python3 build_pptx.py`) so the two stay in sync. Opens
natively in PowerPoint/Keynote, and in Google Slides via Drive's
"Open with Google Slides" on an uploaded .pptx.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------- palette (mirrors the HTML deck's CSS custom properties) ----------
VOID = RGBColor(0x0A, 0x0B, 0x0F)
PANEL = RGBColor(0x14, 0x16, 0x1D)
HAIR = RGBColor(0x2A, 0x28, 0x38)
INK = RGBColor(0xEA, 0xE8, 0xF2)
MUTED = RGBColor(0x8D, 0x90, 0xA6)
FAINT = RGBColor(0x5B, 0x5E, 0x70)
ACCENT = RGBColor(0x9D, 0x7F, 0xFF)
ACCENT_2 = RGBColor(0x6A, 0x48, 0xF0)
GOOD = RGBColor(0x57, 0xD9, 0xA3)
GOOD_SOFT = RGBColor(0x16, 0x2A, 0x24)
WARN = RGBColor(0xFF, 0xB4, 0x54)
WARN_SOFT = RGBColor(0x2C, 0x22, 0x14)

SERIF = "Georgia"
SANS = "Calibri"
MONO = "Consolas"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


def new_slide():
    slide = prs.slides.add_slide(BLANK)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = VOID
    # top accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Pt(4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    bar.shadow.inherit = False
    return slide


def add_text(slide, left, top, width, height, runs, size=18, font=SANS, color=INK,
             bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.0,
             space_after=0):
    """runs: string, or list of (text, {overrides}) tuples for mixed formatting."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    p.space_after = Pt(space_after)
    if isinstance(runs, str):
        runs = [(runs, {})]
    for text, overrides in runs:
        r = p.add_run()
        r.text = text
        r.font.size = Pt(overrides.get("size", size))
        r.font.name = overrides.get("font", font)
        r.font.color.rgb = overrides.get("color", color)
        r.font.bold = overrides.get("bold", bold)
        r.font.italic = overrides.get("italic", False)
    return box


def add_para(tf, runs, size=14, font=SANS, color=MUTED, bold=False, align=PP_ALIGN.LEFT,
             line_spacing=1.15, space_after=6, space_before=0):
    p = tf.add_paragraph() if tf.paragraphs[0].runs or tf.paragraphs[0].text else tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    if isinstance(runs, str):
        runs = [(runs, {})]
    for text, overrides in runs:
        r = p.add_run()
        r.text = text
        r.font.size = Pt(overrides.get("size", size))
        r.font.name = overrides.get("font", font)
        r.font.color.rgb = overrides.get("color", color)
        r.font.bold = overrides.get("bold", bold)
        r.font.italic = overrides.get("italic", False)
    return p


def kicker(slide, text, color=ACCENT, top=Inches(0.55)):
    add_text(slide, Inches(0.8), top, Inches(10), Inches(0.4),
              [("●  ", {"color": color, "font": MONO, "size": 12}),
               (text.upper(), {"color": color, "font": MONO, "size": 12})])


def headline(slide, runs, top=Inches(1.05), size=34, width=Inches(10.5), height=Inches(1.8)):
    add_text(slide, Inches(0.8), top, width, height, runs, size=size, font=SERIF, color=INK,
              line_spacing=1.08)


def chip(slide, left, top, text, kind="good"):
    fg, bg = (GOOD, GOOD_SOFT) if kind == "good" else (WARN, WARN_SOFT)
    w = Inches(0.16 * len(text) + 0.35)
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, Inches(0.32))
    box.adjustments[0] = 0.5
    box.fill.solid(); box.fill.fore_color.rgb = bg
    box.line.color.rgb = fg; box.line.width = Pt(0.75)
    box.shadow.inherit = False
    tf = box.text_frame
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text.upper()
    r.font.size = Pt(9); r.font.name = MONO; r.font.color.rgb = fg
    return box


def card(slide, left, top, width, height, label, desc, tag=None, label_size=13, desc_size=11.5):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.adjustments[0] = 0.045
    box.fill.solid(); box.fill.fore_color.rgb = PANEL
    box.line.color.rgb = HAIR; box.line.width = Pt(0.75)
    box.shadow.inherit = False
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.14); tf.margin_right = Inches(0.14)
    tf.margin_top = Inches(0.1); tf.margin_bottom = Inches(0.1)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = label
    r.font.size = Pt(label_size); r.font.name = MONO; r.font.bold = True; r.font.color.rgb = INK
    if tag:
        r2 = p.add_run(); r2.text = "   " + tag.upper()
        r2.font.size = Pt(8); r2.font.name = MONO; r2.font.color.rgb = GOOD
    p2 = tf.add_paragraph()
    p2.space_before = Pt(4)
    r3 = p2.add_run(); r3.text = desc
    r3.font.size = Pt(desc_size); r3.font.name = SANS; r3.font.color.rgb = MUTED
    return box


def footer(slide, items):
    add_text(slide, Inches(0.8), Inches(7.02), Inches(11), Inches(0.35),
              [(it + ("     " if i < len(items) - 1 else ""), {"color": FAINT, "font": MONO, "size": 10})
               for i, it in enumerate(items)])


# ============================================================ SLIDE 1 — TITLE
s = new_slide()
mark = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.55), Inches(0.9), Inches(0.9))
mark.adjustments[0] = 0.22
mark.fill.solid(); mark.fill.fore_color.rgb = RGBColor(0x15, 0x18, 0x1F)
mark.line.color.rgb = ACCENT; mark.line.width = Pt(1)
mark.shadow.inherit = False
star = s.shapes.add_shape(MSO_SHAPE.STAR_8_POINT, Inches(0.98), Inches(1.73), Inches(0.54), Inches(0.54))
star.fill.solid(); star.fill.fore_color.rgb = ACCENT
star.line.fill.background()
star.shadow.inherit = False
kicker(s, "Product pitch — Anthropic", top=Inches(1.0))
headline(s, [("LakshX", {"size": 54}), ("  IDE", {"size": 16, "color": MUTED, "font": MONO})],
         top=Inches(1.55), size=54, height=Inches(1.1), width=Inches(9))
add_text(s, Inches(1.9), Inches(2.65), Inches(7.2), Inches(1.2),
         "The agentic IDE that remembers — built for one developer today, an entire engineering team tomorrow.",
         size=19, font=SERIF, color=INK, line_spacing=1.25)
chip(s, Inches(0.8), Inches(4.1), "Prototype — shipped", "good")
chip(s, Inches(2.9), Inches(4.1), "Enterprise Agent Brain — ideation", "warn")
footer(s, ["lakshx.in", "MIT licensed", "VS Code fork"])

# ============================================================ SLIDE 2 — PROBLEM
s = new_slide()
kicker(s, "The problem")
headline(s, [("Every AI coding agent today is a ", {}), ("brilliant amnesiac.", {"italic": True, "color": ACCENT})])
tf = s.shapes.add_textbox(Inches(0.8), Inches(2.6), Inches(11), Inches(3)).text_frame
tf.word_wrap = True
rows = [
    ("SESSION-SCOPED", "Close the chat and the reasoning is gone. Tomorrow's session starts from zero, even on the same file."),
    ("MACHINE-SCOPED", "The agent only knows the one folder open on the one laptop it happens to be running on."),
    ("UNVERIFIED “DONE”", "The model tells you it finished. Nothing independently re-checks that claim against reality."),
]
first = True
for tag, txt in rows:
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    first = False
    p.space_after = Pt(14)
    r = p.add_run(); r.text = tag + "   "
    r.font.size = Pt(12); r.font.name = MONO; r.font.color.rgb = ACCENT
    r2 = p.add_run(); r2.text = txt
    r2.font.size = Pt(15); r2.font.name = SANS; r2.font.color.rgb = INK
add_text(s, Inches(0.8), Inches(5.9), Inches(9), Inches(0.6),
         "None of this is a model problem. It's an infrastructure problem.",
         size=17, font=SERIF, color=INK, bold=False)

# ============================================================ SLIDE 3 — INSIGHT
s = new_slide()
kicker(s, "Our insight")
headline(s, [("The next unlock isn't a smarter model. It's ", {}), ("trustworthy autonomy.", {"italic": True, "color": ACCENT})],
         size=32)
cols = [
    ("Verification, not vibes", "“Done” is a re-run of a real, hash-frozen spec — a process exits 0, a test actually passes — not a sentence the model generated."),
    ("Persistent memory, not sessions", "Every action the agent takes is traced and stored locally, always on — the record outlives the chat window."),
    ("Shared context, not one machine", "The long-term bet: that memory shouldn't stop at one person's laptop when the work never did."),
]
cw = Inches(3.7)
for i, (label, desc) in enumerate(cols):
    card(s, Inches(0.8 + i * 3.85), Inches(2.7), cw, Inches(2.4), label, desc, label_size=14, desc_size=12.5)

# ============================================================ SLIDE 4 — WHAT WE'VE BUILT
s = new_slide()
kicker(s, "Product today — shipped")
headline(s, "A VS Code fork with a real engineer inside it.", size=30, height=Inches(1))
feats = [
    ("Royal Mode 2.0", "Plans, executes, and verifies its own work through a gated phase machine."),
    ("Background subagents", "Dispatches and monitors parallel work without blocking the main session."),
    ("Local trace store", "Every agent action recorded locally and always on."),
    ("Bring-your-own-model", "Plug in any provider — nothing locked to one vendor's API."),
    ("Structural search & secrets scan", "AST-aware find/replace and SAST-lite checks built into the editor."),
    ("Dependency & call graph", "A guided tour of any unfamiliar codebase, generated automatically."),
    ("Voice mode", "Offline, local speech-to-text push-to-talk dictation."),
    ("Remote access", "Pair a phone over WiFi to watch and steer a running session."),
]
cw, ch = Inches(2.98), Inches(1.85)
for i, (label, desc) in enumerate(feats):
    col = i % 4; row = i // 4
    card(s, Inches(0.8 + col * 3.08), Inches(2.15 + row * 2.0), cw, ch, label, desc,
         tag="shipped", label_size=12, desc_size=10.5)

# ============================================================ SLIDE 5 — ROYAL MODE DIAGRAM
s = new_slide()
kicker(s, "How it earns trust")
headline(s, [("The agent can't just ", {}), ("say", {"italic": True, "color": ACCENT}), (" it's done.", {})], size=32)

phase_nodes = ["INTAKE", "RECON + PLAN", "EXECUTE", "VERIFY", "DONE"]
x = Inches(0.8); y = Inches(2.85); h = Inches(0.55)
widths = [Inches(1.2), Inches(1.7), Inches(1.3), Inches(1.2), Inches(1.1)]
for i, label in enumerate(phase_nodes):
    w = widths[i]
    is_done = label == "DONE"
    is_verify = label == "VERIFY"
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.adjustments[0] = 0.15
    box.fill.solid()
    box.fill.fore_color.rgb = GOOD_SOFT if is_done else PANEL
    box.line.color.rgb = GOOD if is_done else (ACCENT if is_verify else HAIR)
    box.line.width = Pt(1)
    box.shadow.inherit = False
    tf = box.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    r.font.size = Pt(11); r.font.name = MONO; r.font.color.rgb = GOOD if is_done else INK
    x = x + w + Inches(0.35)
    if i < len(phase_nodes) - 1:
        add_text(s, x - Inches(0.32), y, Inches(0.3), h, "→", size=16, font=SANS, color=FAINT,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

loop1 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(3.75), Inches(5.4), Inches(0.55))
loop1.adjustments[0] = 0.15
loop1.fill.solid(); loop1.fill.fore_color.rgb = WARN_SOFT
loop1.line.color.rgb = WARN; loop1.line.width = Pt(0.75)
loop1.shadow.inherit = False
tf = loop1.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.1)
p = tf.paragraphs[0]
r = p.add_run(); r.text = "VERIFY fails → FIX (≤ 2 rounds) → back to EXECUTE"
r.font.size = Pt(10.5); r.font.name = MONO; r.font.color.rgb = WARN

loop2 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.4), Inches(3.75), Inches(6.0), Inches(0.55))
loop2.adjustments[0] = 0.15
loop2.fill.solid(); loop2.fill.fore_color.rgb = WARN_SOFT
loop2.line.color.rgb = WARN; loop2.line.width = Pt(0.75)
loop2.shadow.inherit = False
tf = loop2.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.1)
p = tf.paragraphs[0]
r = p.add_run(); r.text = "Still wrong → REWIND (≤ 2 re-entries) → back to PLAN"
r.font.size = Pt(10.5); r.font.name = MONO; r.font.color.rgb = WARN

add_text(s, Inches(0.8), Inches(4.7), Inches(10.5), Inches(1.2),
         "VERIFY re-runs a real, hash-frozen spec set before the plan was even executed — the model cannot quietly "
         "redefine “done” partway through a task. If it can't pass, it doesn't get to stop.",
         size=13.5, font=SANS, color=MUTED, line_spacing=1.3)

# ============================================================ SLIDE 6 — COMPETITIVE TABLE
s = new_slide()
kicker(s, "Why not just Cursor / Copilot / Windsurf?")
headline(s, [("They optimize the ", {}), ("typing.", {"italic": True, "color": ACCENT}), (" We're building the ", {}),
             ("trust layer", {"italic": True, "color": ACCENT}), (" underneath it.", {})], size=27, height=Inches(1.3))

rows = [
    ("", "Typical AI IDE", "LakshX"),
    ("Memory survives the session", "— session only", "✓ persistent, local"),
    ("Independently verifies “done”", "— self-reported", "✓ hash-frozen spec re-run"),
    ("Runs background multi-agent work", "— single-threaded chat", "✓ dispatch + monitor"),
    ("Shares context across teammates", "— isolated per machine", "→ Agent Brain (roadmap)"),
    ("Model-agnostic, local-first", "— one vendor, cloud-only", "✓ bring your own model"),
]
gtable = s.shapes.add_table(len(rows), 3, Inches(0.8), Inches(2.85), Inches(11.7), Inches(3.6)).table
gtable.columns[0].width = Inches(4.7)
gtable.columns[1].width = Inches(3.5)
gtable.columns[2].width = Inches(3.5)
for ri, (label, typ, lx) in enumerate(rows):
    is_head = ri == 0
    for ci, text in enumerate((label, typ, lx)):
        cell = gtable.cell(ri, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x1D, 0x16, 0x33) if (ci == 2 and not is_head) else VOID
        cell.margin_left = Inches(0.12); cell.margin_top = Inches(0.06); cell.margin_bottom = Inches(0.06)
        tf = cell.text_frame
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = text
        if is_head:
            r.font.size = Pt(11); r.font.name = MONO; r.font.color.rgb = FAINT
        elif ci == 0:
            r.font.size = Pt(13); r.font.name = SANS; r.font.color.rgb = INK
        elif ci == 1:
            r.font.size = Pt(12); r.font.name = MONO; r.font.color.rgb = FAINT
        else:
            r.font.size = Pt(12); r.font.name = MONO; r.font.color.rgb = GOOD
# strip default table style borders/banding by leaving python-pptx defaults (acceptable for a fast, editable export)

# ============================================================ SLIDE 7 — THE BIG SWING
s = new_slide()
kicker(s, "The big swing", color=WARN)
headline(s, "Every agent today is isolated — to one employee, on one machine, inside one project folder.",
         size=36, top=Inches(1.7), height=Inches(2.6), width=Inches(10.8))
add_text(s, Inches(0.8), Inches(4.6), Inches(10), Inches(1.0),
         "Real engineering isn't done by one person. Why is the agent?",
         size=22, font=SERIF, color=MUTED)

# ============================================================ SLIDE 8 — COORDINATION EXAMPLE
s = new_slide()
kicker(s, "A client project, five engineers", color=WARN)
headline(s, "One team, five modules, five agents that have never met.", size=26, height=Inches(1.0))
rows = [
    ("AUTH DEV", "renames a token field — their agent knows, nobody else's does."),
    ("API DEV", "ships an integration against the old field name — breaks silently in staging, three days later."),
    ("QA ENGINEER", "files the bug. Her agent starts debugging from zero — no idea Tuesday's change happened."),
    ("NEW HIRE", "joins mid-project, re-briefed by humans on decisions five agents already individually “knew.”"),
]
tf = s.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(7.0), Inches(4.4)).text_frame
tf.word_wrap = True
first = True
for tag, txt in rows:
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    first = False
    p.space_after = Pt(16)
    r = p.add_run(); r.text = tag + "  "
    r.font.size = Pt(11); r.font.name = MONO; r.font.color.rgb = ACCENT
    r2 = p.add_run(); r2.text = txt
    r2.font.size = Pt(13.5); r2.font.name = SANS; r2.font.color.rgb = INK

# mini brain diagram, right side
labels = ["auth · agent", "api · agent", "frontend · agent", "infra · agent", "qa · agent"]
by = Inches(2.3)
for i, lab in enumerate(labels):
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.1), by + Inches(i * 0.62), Inches(1.55), Inches(0.42))
    box.adjustments[0] = 0.2
    box.fill.solid(); box.fill.fore_color.rgb = PANEL
    box.line.color.rgb = HAIR; box.line.width = Pt(0.75)
    box.shadow.inherit = False
    tf2 = box.text_frame; tf2.word_wrap = False
    p = tf2.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = lab
    r.font.size = Pt(9.5); r.font.name = MONO; r.font.color.rgb = MUTED
add_text(s, Inches(9.85), by + Inches(1.3), Inches(0.6), Inches(0.5), "→", size=20, color=ACCENT,
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
brain = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.6), Inches(2.85), Inches(1.5), Inches(1.5))
brain.fill.solid(); brain.fill.fore_color.rgb = ACCENT
brain.line.fill.background()
brain.shadow.inherit = False
tf3 = brain.text_frame; tf3.word_wrap = True
p = tf3.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "shared context"
r.font.size = Pt(10.5); r.font.name = MONO; r.font.bold = True; r.font.color.rgb = VOID
add_text(s, Inches(0.8), Inches(6.2), Inches(11.5), Inches(0.9),
         "One shared Agent Brain per project. Any teammate's agent can ask “has anyone touched this, and why?” "
         "and get a real, current answer.",
         size=12.5, font=SANS, color=MUTED, line_spacing=1.25)

# ============================================================ SLIDE 9 — MORE SCENARIOS
s = new_slide()
kicker(s, "This pattern is everywhere", color=WARN)
headline(s, "The five-person team is one example. It's not the only one.", size=27, height=Inches(1.0))
scenarios = [
    ("Agencies with rotating developers", "A contractor rolls off a client project after two months — everything their agent learned evaporates with them. Agent Brain persists at the project level, not the person level."),
    ("2 a.m. incident response", "Whoever's on call starts cold. If this exact failure was already root-caused six weeks ago by a different agent, the Brain says so — and cuts time-to-resolution."),
    ("Regulated industries", "Fintech and healthcare need a provable record of what changed, why, and how it was verified. The verification + trace infrastructure already produces exactly that."),
    ("Distributed, cross-timezone teams", "A handoff today is a Slack message or a stale doc. With a shared Brain, the next timezone's agent already has full context on what happened eight hours ago."),
]
for i, (title, desc) in enumerate(scenarios):
    col = i % 2; row = i // 2
    left = Inches(0.8 + col * 6.0)
    top = Inches(2.2 + row * 2.4)
    bar2 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Pt(2.5), Inches(2.0))
    bar2.fill.solid(); bar2.fill.fore_color.rgb = ACCENT
    bar2.line.fill.background(); bar2.shadow.inherit = False
    tf = s.shapes.add_textbox(left + Inches(0.2), top, Inches(5.5), Inches(2.0)).text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title
    r.font.size = Pt(13); r.font.name = MONO; r.font.color.rgb = INK
    p2 = tf.add_paragraph(); p2.space_before = Pt(6)
    r2 = p2.add_run(); r2.text = desc
    r2.font.size = Pt(11.5); r2.font.name = SANS; r2.font.color.rgb = MUTED

# ============================================================ SLIDE 10 — BUSINESS CASE
s = new_slide()
kicker(s, "The business case")
headline(s, "Companies already pay for coordination. We reduce the bill instead of ignoring it.", size=27, height=Inches(1.5))
rows = [
    ("ALREADY BUDGETED", "Standups, sprint planning, code-review context-gathering, onboarding docs, tribal-knowledge risk on attrition — all real, existing line items on an engineering leader's budget."),
    ("COMPLIANCE ANGLE", "Verification + trace infrastructure, already shipped, turns “an AI agent wrote this” from a liability into a searchable, provable audit trail."),
    ("WHO BUYS", "This sells to a VP Eng or CTO for the whole team — not one developer's individual seat. Bigger contract, stickier, less churn."),
]
tf = s.shapes.add_textbox(Inches(0.8), Inches(3.0), Inches(11.5), Inches(4)).text_frame
tf.word_wrap = True
first = True
for tag, txt in rows:
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    first = False
    p.space_after = Pt(16)
    r = p.add_run(); r.text = tag + "   "
    r.font.size = Pt(11.5); r.font.name = MONO; r.font.color.rgb = ACCENT
    r2 = p.add_run(); r2.text = txt
    r2.font.size = Pt(14); r2.font.name = SANS; r2.font.color.rgb = INK

# ============================================================ SLIDE 11 — FROM LOCAL TO SHARED
s = new_slide()
kicker(s, "Not a pivot — an extension")
headline(s, "The hard infrastructure is already built. Agent Brain is what happens when we open it up.", size=25, height=Inches(1.4))
mapping = [
    ("SHIPPED, SOLO", "Local trace store", "TEAM-SCALE", "Shared, project-level trace index"),
    ("SHIPPED, SOLO", "Background subagent registry", "TEAM-SCALE", "Cross-person task registry"),
    ("SHIPPED, SOLO", "Hash-frozen verification spec", "TEAM-SCALE", "Shared, team-wide “definition of done” per module"),
]
y = Inches(2.9)
for k1, v1, k2, v2 in mapping:
    tfL = s.shapes.add_textbox(Inches(0.8), y, Inches(5.0), Inches(1.0)).text_frame
    tfL.word_wrap = True
    p = tfL.paragraphs[0]
    r = p.add_run(); r.text = k1 + "\n"
    r.font.size = Pt(9.5); r.font.name = MONO; r.font.color.rgb = GOOD
    p2 = tfL.add_paragraph()
    r2 = p2.add_run(); r2.text = v1
    r2.font.size = Pt(14); r2.font.name = SANS; r2.font.color.rgb = INK
    add_text(s, Inches(6.0), y + Inches(0.15), Inches(0.6), Inches(0.5), "→", size=18, color=ACCENT,
              align=PP_ALIGN.CENTER)
    tfR = s.shapes.add_textbox(Inches(6.8), y, Inches(5.6), Inches(1.0)).text_frame
    tfR.word_wrap = True
    p = tfR.paragraphs[0]
    r = p.add_run(); r.text = k2 + "\n"
    r.font.size = Pt(9.5); r.font.name = MONO; r.font.color.rgb = WARN
    p2 = tfR.add_paragraph()
    r2 = p2.add_run(); r2.text = v2
    r2.font.size = Pt(14); r2.font.name = SANS; r2.font.color.rgb = INK
    y += Inches(1.35)

# ============================================================ SLIDE 12 — WHY US
s = new_slide()
kicker(s, "Why we can build this")
headline(s, "We didn't start with the vision slide. We started with the infrastructure it needs.", size=28, height=Inches(1.6))
proofs = [
    ("A shipped verification engine", "Not a plan for one — running, tested, hash-frozen specs re-run against real process exit codes today."),
    ("Shipped multi-agent orchestration", "Dispatching and monitoring parallel background work is the exact substrate a multiplayer Brain needs."),
    ("Shipped structured, persistent traces", "The exact substrate shared memory needs — we just haven't opened it past one machine yet."),
]
for i, (label, desc) in enumerate(proofs):
    card(s, Inches(0.8 + i * 3.85), Inches(3.0), Inches(3.7), Inches(2.4), label, desc, label_size=13.5, desc_size=12)

# ============================================================ SLIDE 13 — BUSINESS MODEL
s = new_slide()
kicker(s, "How this makes money")
headline(s, "Two tiers. Value scales with team size, not just seat count.", size=28, height=Inches(1.1))
card(s, Inches(0.8), Inches(2.7), Inches(5.7), Inches(2.9), "Individual",
     "Seat-based. Bring your own model, run locally, keep every feature on this deck. The prototype you can install right now.",
     tag="live today", label_size=17, desc_size=13)
card(s, Inches(6.75), Inches(2.7), Inches(5.7), Inches(2.9), "Team Agent Brain",
     "Org-wide. Priced on team/project, not per head — the more people share one project's Brain, the more value it returns. Natural expansion revenue as teams grow.",
     tag=None, label_size=17, desc_size=13)
chip(s, Inches(9.1), Inches(2.78), "roadmap", "warn")

# ============================================================ SLIDE 14 — STATUS
s = new_slide()
kicker(s, "Where we are")
headline(s, "Prototype. Ideation phase for the enterprise layer. Exactly where we say we are.", size=26, height=Inches(1.5))
stages = ["PROTOTYPE — SHIPPED", "DESIGNING THE SHARED-CONTEXT PROTOCOL", "TEAM PILOT", "GENERAL AVAILABILITY"]
widths2 = [Inches(2.1), Inches(3.5), Inches(1.7), Inches(2.5)]
x = Inches(0.8); y = Inches(3.1)
for i, label in enumerate(stages):
    w = widths2[i]
    is_shipped = i == 0
    is_active = i == 1
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, Inches(0.6))
    box.adjustments[0] = 0.12
    box.fill.solid()
    box.fill.fore_color.rgb = GOOD_SOFT if is_shipped else PANEL
    box.line.color.rgb = GOOD if is_shipped else (ACCENT if is_active else HAIR)
    box.line.width = Pt(1)
    box.shadow.inherit = False
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    r.font.size = Pt(9.5); r.font.name = MONO; r.font.color.rgb = GOOD if is_shipped else INK
    x = x + w + Inches(0.3)
    if i < len(stages) - 1:
        add_text(s, x - Inches(0.28), y, Inches(0.26), Inches(0.6), "→", size=14, color=FAINT,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(0.8), Inches(4.3), Inches(11), Inches(1.0),
         "No overclaiming: the single-player product on slide 4 is real and running. Enterprise Agent Brain is the "
         "ideation-stage bet this raise is for.", size=13, font=SANS, color=MUTED, line_spacing=1.3)

# ============================================================ SLIDE 15 — WHY ANTHROPIC
s = new_slide()
kicker(s, "Built with Claude, built for Claude")
headline(s, "This isn't a “we'll add Claude support” pitch. Claude is the default, top to bottom.",
         size=27, height=Inches(1.5))
why_anthropic = [
    ("The product's own default model", "A fresh LakshX install's out-of-the-box default provider is Anthropic, model Claude Sonnet 5 — not a configuration option buried in settings, the actual shipped default every new user starts from."),
    ("Built substantially with Claude Code itself", "A real, verifiable share of this codebase's own commit history carries “Co-Authored-By: Claude Sonnet 5” — Royal Mode's verification engine, the background subagent registry, and this deck's own bug fixes were built in collaboration with Claude Code, not just marketed as AI-friendly."),
]
for i, (label, desc) in enumerate(why_anthropic):
    card(s, Inches(0.8 + i * 3.85), Inches(2.9), Inches(3.7), Inches(2.6), label, desc, label_size=13, desc_size=11.5)
add_text(s, Inches(0.8), Inches(5.9), Inches(11.5), Inches(1.0),
         "The pitch: LakshX is a genuine showcase of Claude's agentic tool-use ceiling — verification-gated autonomy, "
         "background multi-agent orchestration, structured traces — running as the default experience, not a bolted-on integration.",
         size=12.5, font=SANS, color=MUTED, line_spacing=1.3)

# ============================================================ SLIDE 16 — THE ASK
s = new_slide()
kicker(s, "What we're applying for")
headline(s, "API credits and rate limits to keep building — and a design partner as Royal Mode pushes on Claude's tool-use ceiling.",
         size=22, height=Inches(1.6))
asks = [
    ("API credits + rate limits", "To cover development traffic now, and the shared Agent Brain service's usage as it moves from ideation to a real pilot — both run on Claude by default, so this is directly proportional to Anthropic's own usage, not a side request."),
    ("Early access + a real feedback loop", "Royal Mode's verification-gated autonomy and background subagent orchestration push hard on tool-use and long-horizon agentic behavior — exactly the surface area new Claude capabilities land on first. Happy to be a design partner / early-access tester and report back with real, shipped-product data, not a lab benchmark."),
]
for i, (who, desc) in enumerate(asks):
    card(s, Inches(0.8 + i * 5.85), Inches(3.0), Inches(5.7), Inches(2.6), who, desc, label_size=13, desc_size=11.5)
add_text(s, Inches(0.8), Inches(5.9), Inches(11), Inches(0.8),
         "Raise amount, equity, and cohort details: fill in before submitting — deliberately left open here.",
         size=12, font=SANS, color=FAINT)

# ============================================================ SLIDE 17 — CLOSING
s = new_slide()
kicker(s, "Where this goes")
headline(s, [("From autocomplete. To an autonomous agent. To an engineering team's ", {}),
             ("shared brain.", {"italic": True, "color": ACCENT})],
         size=34, top=Inches(2.2), height=Inches(3), width=Inches(11))
footer(s, ["lakshx.in", "github.com/BeastxD7/LakshX-IDE", "contact@lakshx.in"])

prs.save("lakshx-pitch-deck.pptx")
print("wrote lakshx-pitch-deck.pptx —", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
